import os
import re
import threading
import tkinter as tk
from tkinter import filedialog, scrolledtext, ttk
from pptx import Presentation
import fitz  # PyMuPDF for PDF files
from docx import Document  # For .docx files

def browse_directory():
    directory = filedialog.askdirectory()
    entry_directory.delete(0, tk.END)
    entry_directory.insert(0, directory)

def search_action():
    threading.Thread(target=run_search, daemon=True).start()

def search_files(directory, pattern, file_extension='*', show_lines=False, progress_callback=None):
    results = []
    regex = re.compile(pattern, re.IGNORECASE)  # Case-insensitive search
    file_count = sum(len(files) for _, _, files in os.walk(directory))
    scanned_files = 0

    for root_dir, _, files in os.walk(directory):
        for file in files:
            if file_extension == '*' or file.endswith(file_extension):
                filepath = os.path.join(root_dir, file)
                try:
                    if file.endswith(('.txt', '.log', '.py')):
                        with open(filepath, "r", encoding="utf-8", errors='ignore') as f:
                            for line_number, line in enumerate(f, start=1):
                                if regex.search(line):
                                    results.append(f"{filepath} (Line {line_number}): {line.strip()}")
                    elif file.endswith('.pptx'):
                        ppt = Presentation(filepath)
                        for slide_number, slide in enumerate(ppt.slides, start=1):
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and regex.search(shape.text):
                                    results.append(f"{filepath} (Slide {slide_number}): {shape.text.strip()}")
                    elif file.endswith('.pdf'):
                        pdf_doc = fitz.open(filepath)
                        for page_num, page in enumerate(pdf_doc, start=1):
                            text = page.get_text("text")
                            if regex.search(text):
                                results.append(f"{filepath} (Page {page_num}): {text.strip()}")
                    elif file.endswith('.docx'):
                        doc = Document(filepath)
                        for para_num, para in enumerate(doc.paragraphs, start=1):
                            if regex.search(para.text):
                                results.append(f"{filepath} (Paragraph {para_num}): {para.text.strip()}")
                except Exception as e:
                    results.append(f"Error reading {filepath}: {e}")
                
                scanned_files += 1
                if progress_callback:
                    progress_callback(scanned_files, file_count)
    
    return results

def run_search():
    directory = entry_directory.get()
    pattern = entry_pattern.get()
    file_extension = file_type_var.get()
    show_lines = show_lines_var.get()
    
    if not directory or not pattern:
        text_results.insert(tk.END, "Please select a directory and enter a search pattern.\n")
        return
    
    text_results.delete('1.0', tk.END)
    progress_bar.start()
    
    results = search_files(directory, pattern, file_extension, show_lines, update_progress)
    
    text_results.insert(tk.END, "\n".join(results) if results else "No matches found.\n")
    save_results(results)
    progress_bar.stop()

def update_progress(current, total):
    progress_bar["value"] = (current / total) * 100
    root.update_idletasks()

def save_results(results, output_file="results.txt"):
    with open(output_file, "w", encoding="utf-8") as f:
        f.write("\n".join(results))

# GUI Setup
root = tk.Tk()
root.title("File Search Tool")
root.geometry("600x450")

frame = tk.Frame(root)
frame.pack(pady=10)

tk.Label(frame, text="Directory:").grid(row=0, column=0, padx=5)
entry_directory = tk.Entry(frame, width=50)
entry_directory.grid(row=0, column=1, padx=5)
tk.Button(frame, text="Browse", command=browse_directory).grid(row=0, column=2, padx=5)

tk.Label(frame, text="Search Pattern:").grid(row=1, column=0, padx=5)
entry_pattern = tk.Entry(frame, width=50)
entry_pattern.grid(row=1, column=1, padx=5)

tk.Label(frame, text="File Type:").grid(row=2, column=0, padx=5)
file_type_var = tk.StringVar(value='*')
tk.OptionMenu(frame, file_type_var, '*', '.txt', '.log', '.py', '.pptx', '.pdf', '.docx').grid(row=2, column=1, padx=5)

show_lines_var = tk.BooleanVar(value=True)
tk.Checkbutton(frame, text="Show Matching Lines", variable=show_lines_var).grid(row=3, column=1, padx=5)

tk.Button(frame, text="Search", command=search_action).grid(row=4, column=1, pady=10)

progress_bar = ttk.Progressbar(root, orient=tk.HORIZONTAL, length=500, mode='determinate')
progress_bar.pack(pady=5)

text_results = scrolledtext.ScrolledText(root, width=70, height=15)
text_results.pack(pady=5)

root.mainloop()
